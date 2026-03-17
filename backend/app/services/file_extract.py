from io import BytesIO
from pathlib import Path

from PIL import Image
from pypdf import PdfReader
from pptx import Presentation
import pytesseract

from app.services.text_normalizer import decode_text_bytes, normalize_text


ALLOWED_EXTENSIONS = {
    ".txt",
    ".md",
    ".pdf",
    ".pptx",
    ".png",
    ".jpg",
    ".jpeg",
    ".bmp",
    ".tiff",
}


def _extract_pdf(data: bytes) -> str:
    reader = PdfReader(BytesIO(data))
    parts = [page.extract_text() or "" for page in reader.pages]
    return normalize_text("\n".join(parts))


def _extract_pptx(data: bytes) -> str:
    presentation = Presentation(BytesIO(data))
    texts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text.strip())
    return normalize_text("\n".join(texts))


def _extract_image(data: bytes) -> str:
    image = Image.open(BytesIO(data))
    return normalize_text(pytesseract.image_to_string(image))


def extract_text_from_file(filename: str, data: bytes) -> str:
    ext = Path(filename).suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise ValueError(f"지원하지 않는 파일 형식입니다: {ext}")

    if ext in {".txt", ".md"}:
        return decode_text_bytes(data)
    if ext == ".pdf":
        return _extract_pdf(data)
    if ext == ".pptx":
        return _extract_pptx(data)
    return _extract_image(data)
